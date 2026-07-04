"""
Generate Django + GesPresence PowerPoint presentation.
Run: .\venv\Scripts\python.exe generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color Palette ─────────────────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x0A, 0x5E, 0xB0)   # #0A5EB0  Django/OFPPT blue
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)   # #E8F0FE  light background
GREEN       = RGBColor(0x19, 0x87, 0x54)   # #198754  success
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # #FF6B35  accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x21, 0x25, 0x29)   # near black
GRAY        = RGBColor(0x6C, 0x75, 0x7D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def header_bar(slide, title, subtitle=None):
    """Blue header bar at top."""
    add_rect(slide, 0, 0, 13.33, 1.4, fill_color=BLUE_DARK)
    add_text(slide, title, 0.4, 0.1, 12, 0.8, size=32, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12, 0.5, size=14,
                 color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)

def bullet_list(slide, items, l, t, w, h, size=15, color=DARK, spacing=0.38):
    """Render a list of (indent, text) tuples as bullet points."""
    y = t
    for (indent, item) in items:
        bullet = '    ' * indent + ('• ' if indent == 0 else '  ◦ ') + item
        add_text(slide, bullet, l, y, w, spacing + 0.05, size=size, color=color)
        y += spacing

def code_box(slide, code, l, t, w, h):
    """Dark code block."""
    add_rect(slide, l, t, w, h, fill_color=RGBColor(0x1E, 0x1E, 0x2E))
    add_rect(slide, l, t, w, 0.3, fill_color=RGBColor(0x31, 0x31, 0x4A))
    add_text(slide, 'python', l + 0.1, t + 0.02, 2, 0.28, size=11,
             color=RGBColor(0xBB, 0xBB, 0xFF), italic=True)
    txb = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.35),
                                    Inches(w - 0.3), Inches(h - 0.45))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = 'Courier New'
        run.font.size  = Pt(11)
        run.font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)

def tag_pill(slide, text, l, t, color=BLUE_DARK):
    add_rect(slide, l, t, len(text)*0.095 + 0.2, 0.32, fill_color=color)
    add_text(slide, text, l + 0.08, t + 0.03, len(text)*0.1 + 0.1, 0.28,
             size=11, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BLUE_DARK)
# diagonal accent
add_rect(s, 0, 4.5, 13.33, 3, fill_color=RGBColor(0x07, 0x46, 0x8A))
# green bar
add_rect(s, 0, 4.4, 13.33, 0.12, fill_color=GREEN)

add_text(s, '🐍 Django Framework', 0.6, 1.0, 12, 1.0,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, 'Construire une Application Web Fullstack', 0.6, 2.1, 11, 0.7,
         size=26, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
add_text(s, 'De zéro à une app complète de gestion de présence', 0.6, 2.85, 11, 0.5,
         size=18, color=RGBColor(0xAA, 0xBB, 0xDD), italic=True, align=PP_ALIGN.LEFT)

add_rect(s, 0.6, 3.5, 0.06, 0.06, fill_color=GREEN)
add_text(s, 'Python • Django 4.2 LTS • MySQL • Bootstrap 5', 0.72, 3.42, 10, 0.4,
         size=14, color=RGBColor(0x88, 0xFF, 0x99))

add_text(s, '2026 — GesPresence Project', 0.6, 6.8, 8, 0.4,
         size=12, color=RGBColor(0x88, 0xAA, 0xCC), italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Plan de la Présentation', 'Ce que vous allez apprendre')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=BLUE_LIGHT)

sections = [
    ('1', 'Introduction à Django',         '2-3',  BLUE_DARK),
    ('2', 'Architecture MVT',              '4-5',  GREEN),
    ('3', 'Modèles & Base de données',     '6-7',  ORANGE),
    ('4', 'Vues & URLs',                   '8-9',  RGBColor(0x6F, 0x42, 0xC1)),
    ('5', 'Templates & Formulaires',       '10-11',BLUE_DARK),
    ('6', 'Authentification & Admin',      '12-13',GREEN),
    ('7', 'Déploiement',                   '14',   ORANGE),
    ('8', '🎓 Projet GesPresence',        '15-22',RGBColor(0xDC, 0x35, 0x45)),
]

for i, (num, title, pages, color) in enumerate(sections):
    col = i % 4
    row = i // 4
    x = 0.35 + col * 3.2
    y = 1.7 + row * 1.8
    add_rect(s, x, y, 2.9, 1.5, fill_color=WHITE)
    add_rect(s, x, y, 2.9, 0.45, fill_color=color)
    add_text(s, num, x + 0.1, y + 0.05, 0.4, 0.38, size=20, bold=True, color=WHITE)
    add_text(s, title, x + 0.12, y + 0.52, 2.7, 0.6, size=13, bold=True, color=DARK)
    add_text(s, f'Diapo {pages}', x + 0.12, y + 1.15, 2, 0.3, size=10, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — QU'EST-CE QUE DJANGO ?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Qu'est-ce que Django ?", 'Le framework web Python "batteries included"')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8, 0xF9, 0xFA))

add_text(s, 'Django est un framework web Python de haut niveau qui encourage\nle développement rapide et une conception propre et pragmatique.', 0.4, 1.55, 8.5, 0.9, size=15, color=DARK)

facts = [
    (0, 'Créé en 2003 par Adrian Holovaty & Simon Willison'),
    (0, 'Open-source depuis 2005'),
    (0, 'Suit le principe "Don\'t Repeat Yourself" (DRY)'),
    (0, 'Inclut : ORM, Admin, Auth, Forms, Templates, Security...'),
    (0, 'Utilisé par : Instagram, Pinterest, Disqus, Mozilla, Spotify'),
]
bullet_list(s, facts, 0.4, 2.55, 8.5, 3.5, size=14)

# Right panel
add_rect(s, 9.1, 1.5, 3.9, 5.7, fill_color=BLUE_DARK)
add_text(s, 'Pourquoi Django ?', 9.3, 1.6, 3.5, 0.5, size=16, bold=True, color=WHITE)
advantages = [
    '⚡ Rapide à développer',
    '🔒 Sécurisé par défaut',
    '📦 Batteries incluses',
    '🌍 Très documenté',
    '🔧 Scalable',
    '👑 Django Admin gratuit',
    '🗄️ ORM puissant',
]
for i, adv in enumerate(advantages):
    add_text(s, adv, 9.3, 2.2 + i * 0.62, 3.5, 0.55, size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE MVT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Architecture MVT', 'Model — View — Template')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF0, 0xF4, 0xFF))

# Three boxes
boxes = [
    ('M', 'Model', 'Définit la structure des\ndonnées et interagit\navec la base de données\nvia l\'ORM Django', BLUE_DARK),
    ('V', 'View', 'Contient la logique\nmétier. Reçoit les\nrequêtes, interroge\nles models et retourne\nune réponse', GREEN),
    ('T', 'Template', 'Fichiers HTML avec\nsyntaxe Django pour\nafficher les données\ndynamiquement', ORANGE),
]
for i, (letter, title, desc, color) in enumerate(boxes):
    x = 0.5 + i * 4.1
    add_rect(s, x, 1.6, 3.7, 5.5, fill_color=WHITE)
    add_rect(s, x, 1.6, 3.7, 0.8, fill_color=color)
    add_text(s, letter, x + 0.15, 1.65, 0.6, 0.7, size=36, bold=True, color=WHITE)
    add_text(s, title, x + 0.75, 1.72, 2.8, 0.6, size=22, bold=True, color=WHITE)
    add_text(s, desc, x + 0.2, 2.55, 3.3, 3.5, size=13, color=DARK)

# Arrows
add_text(s, '→', 4.25, 4.1, 0.5, 0.5, size=24, bold=True, color=BLUE_DARK)
add_text(s, '→', 8.35, 4.1, 0.5, 0.5, size=24, bold=True, color=GREEN)

# Browser & DB
add_rect(s, 0.0, 6.7, 2.5, 0.55, fill_color=RGBColor(0xDC, 0x35, 0x45))
add_text(s, '🌐 Navigateur', 0.1, 6.72, 2.3, 0.45, size=13, bold=True, color=WHITE)
add_rect(s, 10.8, 6.7, 2.5, 0.55, fill_color=RGBColor(0x6F, 0x42, 0xC1))
add_text(s, '🗄️ Base de données', 10.85, 6.72, 2.4, 0.45, size=13, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FLUX D'UNE REQUÊTE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Flux d\'une Requête HTTP', 'Comment Django traite une requête')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8, 0xF9, 0xFA))

steps = [
    ('1', 'Navigateur',    'GET /etudiants/', RGBColor(0xDC,0x35,0x45)),
    ('2', 'urls.py',       'path(\'etudiants/\',\nviews.liste)', BLUE_DARK),
    ('3', 'views.py',      'def liste(request):\n  data = Model.objects\n        .all()', GREEN),
    ('4', 'models.py',     'SELECT * FROM\netudiants;', RGBColor(0x6F,0x42,0xC1)),
    ('5', 'template.html', '{% for e in data %}\n  {{ e.nom }}\n{% endfor %}', ORANGE),
    ('6', 'Response',      'HTML → Navigateur', RGBColor(0x0D,0xCA,0xF0)),
]

for i, (num, title, code_text, color) in enumerate(steps):
    x = 0.2 + i * 2.15
    add_rect(s, x, 1.55, 2.0, 0.5, fill_color=color)
    add_text(s, f'{num}. {title}', x + 0.08, 1.6, 1.9, 0.42, size=13, bold=True, color=WHITE)
    add_rect(s, x, 2.1, 2.0, 2.0, fill_color=RGBColor(0x1E,0x1E,0x2E))
    add_text(s, code_text, x + 0.1, 2.18, 1.85, 1.85, size=10,
             color=RGBColor(0xCC,0xFF,0xCC))
    if i < 5:
        add_text(s, '→', x + 2.02, 2.4, 0.2, 0.4, size=20, bold=True, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODÈLES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Les Modèles Django', 'ORM — Object Relational Mapper')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

add_text(s, 'Un modèle = une table en base de données', 0.4, 1.55, 8, 0.45, size=16, bold=True, color=BLUE_DARK)

code1 = """# models.py
from django.db import models

class Etudiant(models.Model):
    nom    = models.CharField(max_length=100)
    email  = models.EmailField(unique=True)
    groupe = models.ForeignKey(
                 'Groupe',
                 on_delete=models.CASCADE
             )
    photo  = models.ImageField(
                 upload_to='etudiants/',
                 blank=True
             )

    def __str__(self):
        return self.nom

    class Meta:
        ordering = ['nom']"""
code_box(s, code1, 0.3, 2.05, 5.8, 5.1)

right_items = [
    (0, 'Types de champs courants :'),
    (1, 'CharField, TextField'),
    (1, 'IntegerField, FloatField'),
    (1, 'DateField, DateTimeField'),
    (1, 'BooleanField, ImageField'),
    (1, 'EmailField, URLField'),
    (0, 'Relations :'),
    (1, 'ForeignKey (Many-to-One)'),
    (1, 'ManyToManyField'),
    (1, 'OneToOneField'),
    (0, 'Commandes migrations :'),
    (1, 'python manage.py makemigrations'),
    (1, 'python manage.py migrate'),
]
bullet_list(s, right_items, 6.4, 1.9, 6.5, 5, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ORM QUERIES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'ORM Django — Requêtes', 'Interroger la base de données sans SQL')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

queries = [
    ('Tous les objets', 'Etudiant.objects.all()'),
    ('Filtrer', 'Etudiant.objects.filter(groupe__nom="DEV-101")'),
    ('Un objet', 'Etudiant.objects.get(id=1)'),
    ('Exclure', 'Etudiant.objects.exclude(actif=False)'),
    ('Compter', 'Etudiant.objects.filter(absent=True).count()'),
    ('Trier', 'Etudiant.objects.order_by("-date_inscription")'),
    ('Créer', 'Etudiant.objects.create(nom="Ahmed", email="a@b.com")'),
    ('Modifier', 'e = Etudiant.objects.get(id=1)\ne.nom = "Ali"\ne.save()'),
    ('Supprimer', 'Etudiant.objects.filter(actif=False).delete()'),
]
for i, (label, code) in enumerate(queries):
    row = i % 5
    col = i // 5
    x = 0.3 + col * 6.5
    y = 1.7 + row * 1.08
    add_rect(s, x, y, 6.2, 0.95, fill_color=WHITE)
    add_rect(s, x, y, 1.8, 0.95, fill_color=BLUE_DARK)
    add_text(s, label, x + 0.08, y + 0.22, 1.7, 0.5, size=12, bold=True, color=WHITE)
    add_text(s, code, x + 1.9, y + 0.08, 4.2, 0.8, size=11,
             color=RGBColor(0x19,0x87,0x54))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — VUES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Les Vues (Views)', 'La logique métier de votre application')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

add_text(s, 'Vues basées sur des fonctions (FBV)', 0.4, 1.6, 6, 0.4, size=15, bold=True, color=BLUE_DARK)
code2 = """from django.shortcuts import render, get_object_or_404
from django.contrib.auth.decorators import login_required
from .models import Etudiant

@login_required
def liste_etudiants(request):
    etudiants = Etudiant.objects.all()
    return render(request,
        'etudiants/liste.html',
        {'etudiants': etudiants}
    )

def detail_etudiant(request, pk):
    etudiant = get_object_or_404(Etudiant, pk=pk)
    return render(request,
        'etudiants/detail.html',
        {'etudiant': etudiant}
    )"""
code_box(s, code2, 0.3, 2.05, 6.1, 5.1)

add_text(s, 'URLs — Routage', 6.7, 1.6, 6.3, 0.4, size=15, bold=True, color=GREEN)
code3 = """# urls.py
from django.urls import path
from . import views

urlpatterns = [
    path('',
         views.liste_etudiants,
         name='etudiant_list'),

    path('<int:pk>/',
         views.detail_etudiant,
         name='etudiant_detail'),
]"""
code_box(s, code3, 6.7, 2.05, 6.3, 4.0)

# Decorators note
add_rect(s, 6.7, 6.2, 6.3, 0.9, fill_color=RGBColor(0xFF,0xF3,0xCD))
add_text(s, '⚠️ Décorateurs utiles: @login_required  @permission_required  @staff_member_required',
         6.85, 6.28, 6.0, 0.6, size=12, color=RGBColor(0x85, 0x6A, 0x04))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TEMPLATES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Les Templates Django', 'Afficher les données en HTML')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

code4 = """<!-- templates/etudiants/liste.html -->
{% extends 'base.html' %}
{% block title %}Liste Étudiants{% endblock %}

{% block content %}
<h2>{{ groupe.nom }} — {{ etudiants.count }} étudiant(s)</h2>

<table class="table">
  {% for etudiant in etudiants %}
  <tr>
    <td>{{ etudiant.nom }}</td>
    <td>{{ etudiant.email }}</td>
    <td>
      {% if etudiant.taux >= 75 %}
        <span class="badge bg-success">✓ OK</span>
      {% else %}
        <span class="badge bg-danger">⚠ Risque</span>
      {% endif %}
    </td>
    <td>
      <a href="{% url 'etudiant_detail' etudiant.pk %}">
        Voir
      </a>
    </td>
  </tr>
  {% empty %}
    <tr><td>Aucun étudiant.</td></tr>
  {% endfor %}
</table>
{% endblock %}"""
code_box(s, code4, 0.3, 1.6, 6.3, 5.7)

tags = [
    ('Tags de contrôle', [
        '{% if %} {% elif %} {% else %} {% endif %}',
        '{% for x in liste %} {% endfor %}',
        '{% block nom %} {% endblock %}',
        '{% extends "base.html" %}',
        '{% include "partial.html" %}',
        '{% url "nom_vue" pk %}',
        '{% csrf_token %}',
        '{% load tags_perso %}',
    ]),
    ('Filtres', [
        '{{ texte|upper }}',
        '{{ date|date:"d/m/Y" }}',
        '{{ liste|length }}',
        '{{ valeur|default:"N/A" }}',
        '{{ nb|floatformat:1 }}%',
    ]),
]
y = 1.65
for (title, items) in tags:
    add_text(s, title, 6.8, y, 6.2, 0.38, size=14, bold=True, color=BLUE_DARK)
    y += 0.38
    for item in items:
        add_rect(s, 6.8, y, 6.2, 0.35, fill_color=RGBColor(0x1E,0x1E,0x2E))
        add_text(s, item, 6.9, y + 0.04, 6.0, 0.3, size=11,
                 color=RGBColor(0xCC,0xFF,0xCC))
        y += 0.38
    y += 0.15

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FORMULAIRES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Les Formulaires Django', 'Validation et traitement des données')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

code5 = """# forms.py
from django import forms
from .models import Etudiant

class EtudiantForm(forms.ModelForm):
    class Meta:
        model  = Etudiant
        fields = ['nom', 'email', 'groupe']
        widgets = {
            'nom': forms.TextInput(
                attrs={'class': 'form-control'}
            ),
        }

# views.py
def creer_etudiant(request):
    if request.method == 'POST':
        form = EtudiantForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('liste')
    else:
        form = EtudiantForm()
    return render(request,
        'form.html', {'form': form})"""
code_box(s, code5, 0.3, 1.6, 5.8, 5.6)

code6 = """<!-- templates/form.html -->
<form method="post">
  {% csrf_token %}
  {{ form.as_p }}
  <button type="submit"
    class="btn btn-primary">
    Enregistrer
  </button>
</form>

<!-- Ou champ par champ : -->
<form method="post">
  {% csrf_token %}
  {{ form.nom.label_tag }}
  {{ form.nom }}
  {{ form.nom.errors }}
  <button type="submit">OK</button>
</form>"""
code_box(s, code6, 6.3, 1.6, 5.8, 4.2)

types_forms = [
    (0, 'ModelForm — basé sur un modèle'),
    (0, 'Form — formulaire personnalisé'),
    (0, 'Validation automatique des champs'),
    (0, 'Protection CSRF intégrée'),
    (0, 'django-crispy-forms pour Bootstrap'),
]
bullet_list(s, types_forms, 6.3, 5.9, 6, 1.5, size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — AUTHENTIFICATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Authentification Django', 'Système d\'auth complet intégré')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

code7 = """# Modèle utilisateur personnalisé
from django.contrib.auth.models import (
    AbstractUser
)

class Utilisateur(AbstractUser):
    ROLES = [
        ('admin',      'Administrateur'),
        ('enseignant', 'Enseignant'),
        ('etudiant',   'Étudiant'),
        ('parent',     'Parent'),
    ]
    role      = models.CharField(
                    max_length=20,
                    choices=ROLES
                )
    telephone = models.CharField(
                    max_length=20,
                    blank=True
                )

# settings.py
AUTH_USER_MODEL = 'accounts.Utilisateur'"""
code_box(s, code7, 0.3, 1.6, 5.8, 5.6)

auth_items = [
    (0, 'Authentification intégrée :'),
    (1, 'django.contrib.auth'),
    (1, 'login() / logout() / authenticate()'),
    (1, '@login_required decorator'),
    (0, 'Vues prêtes à l\'emploi :'),
    (1, 'LoginView, LogoutView'),
    (1, 'PasswordChangeView'),
    (1, 'PasswordResetView'),
    (0, 'Permissions & Groupes :'),
    (1, 'user.has_perm("app.permission")'),
    (1, '@permission_required("app.perm")'),
    (0, 'Session management automatique'),
    (0, 'Protection contre CSRF, XSS, SQLi'),
]
bullet_list(s, auth_items, 6.4, 1.65, 6.5, 5.5, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DJANGO ADMIN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Django Admin', 'Interface d\'administration automatique et gratuite')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

code8 = """# admin.py
from django.contrib import admin
from .models import Etudiant, Groupe

@admin.register(Etudiant)
class EtudiantAdmin(admin.ModelAdmin):
    list_display = (
        'matricule',
        'nom',
        'groupe',
        'taux_presence'
    )
    list_filter  = ('groupe__filiere',)
    search_fields = ('matricule', 'nom')
    date_hierarchy = 'date_inscription'

    def taux_presence(self, obj):
        return f"{obj.get_taux()}%"
    taux_presence.short_description = 'Taux'

admin.site.register(Groupe)"""
code_box(s, code8, 0.3, 1.6, 5.8, 5.5)

features = [
    '✅ CRUD automatique pour tous les modèles',
    '✅ Recherche, filtres, tri intégrés',
    '✅ Gestion des relations (inline)',
    '✅ Export, import de données',
    '✅ Permissions par modèle et par utilisateur',
    '✅ Historique des modifications',
    '✅ Interface responsive Bootstrap',
    '',
    '🌐 Accès : /admin/',
    '👤 Création : python manage.py createsuperuser',
]
y = 1.65
for feat in features:
    add_text(s, feat, 6.4, y, 6.5, 0.42, size=13, color=DARK if feat else GRAY)
    y += 0.48

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — STATIC FILES & DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Fichiers Statiques & Déploiement', 'CSS, JS, Images et mise en production')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

code9 = """# settings.py
STATIC_URL  = '/static/'
STATICFILES_DIRS = [BASE_DIR / 'static']
STATIC_ROOT = BASE_DIR / 'staticfiles'

MEDIA_URL   = '/media/'
MEDIA_ROOT  = BASE_DIR / 'media'

# Template
{% load static %}
<link href="{% static 'css/main.css' %}"
      rel="stylesheet">
<img src="{% static 'img/logo.png' %}">

# Collect (production)
$ python manage.py collectstatic"""
code_box(s, code9, 0.3, 1.6, 5.8, 4.5)

deploy_steps = [
    ('1', 'requirements.txt', 'pip freeze > requirements.txt', BLUE_DARK),
    ('2', 'Variables .env',   'SECRET_KEY, DEBUG=False, DB config', GREEN),
    ('3', 'collectstatic',    'python manage.py collectstatic', ORANGE),
    ('4', 'Gunicorn',         'gunicorn gespresence.wsgi:application', RGBColor(0x6F,0x42,0xC1)),
    ('5', 'Nginx / Apache',   'Reverse proxy + SSL/HTTPS', RGBColor(0xDC,0x35,0x45)),
]
for i, (n, title, cmd, color) in enumerate(deploy_steps):
    y = 1.65 + i * 1.0
    add_rect(s, 6.4, y, 6.5, 0.85, fill_color=WHITE)
    add_rect(s, 6.4, y, 0.55, 0.85, fill_color=color)
    add_text(s, n, 6.5, y + 0.18, 0.4, 0.5, size=16, bold=True, color=WHITE)
    add_text(s, title, 7.05, y + 0.05, 5.6, 0.35, size=13, bold=True, color=DARK)
    add_text(s, cmd, 7.05, y + 0.45, 5.6, 0.35, size=11, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — TRANSITION PROJET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=DARK)
add_rect(s, 0, 3.2, 13.33, 0.12, fill_color=RGBColor(0xDC,0x35,0x45))

add_text(s, 'Projet Pratique', 0.6, 0.8, 12, 0.8,
         size=22, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, '🎓 GesPresence', 0.6, 1.6, 12, 1.2,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 'Application Web de Gestion de Présence des Étudiants',
         0.6, 2.85, 12, 0.6,
         size=20, color=RGBColor(0xAA,0xBB,0xFF), align=PP_ALIGN.CENTER)
add_text(s, 'Django 4.2  •  MySQL  •  Bootstrap 5  •  Python 3.13',
         0.6, 3.5, 12, 0.5,
         size=15, color=RGBColor(0x88,0xFF,0x99), align=PP_ALIGN.CENTER,
         italic=True)

tags_list = [
    (2.5,  5.0, 'Multi-rôles', BLUE_DARK),
    (5.0,  5.0, 'Présences',   GREEN),
    (7.3,  5.0, 'Rapports',    ORANGE),
    (9.4,  5.0, 'Notifications', RGBColor(0x6F,0x42,0xC1)),
    (3.5,  5.7, 'Export Excel', RGBColor(0xDC,0x35,0x45)),
    (6.0,  5.7, 'Admin Django', BLUE_DARK),
    (8.3,  5.7, 'ORM MySQL',   GREEN),
]
for (lx, ty, label, color) in tags_list:
    tag_pill(s, label, lx, ty, color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONTEXTE PROJET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'GesPresence — Contexte & Objectifs', 'Application de gestion de présence OFPPT')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

add_text(s, 'Problématique', 0.4, 1.55, 12, 0.42, size=18, bold=True, color=BLUE_DARK)
add_text(s,
    'Les établissements OFPPT avaient besoin d\'un système numérique pour remplacer les\n'
    'feuilles de présence papier et permettre un suivi en temps réel des absences.',
    0.4, 2.0, 12.5, 0.75, size=14, color=DARK)

cols = [
    ('🎯 Objectifs', [
        'Digitaliser le marquage des présences',
        'Alerter les parents en cas d\'absence',
        'Générer des rapports statistiques',
        'Gérer filières, groupes et emplois du temps',
    ], BLUE_DARK),
    ('👥 Utilisateurs', [
        '👤 Administrateur — gestion complète',
        '👨‍🏫 Enseignant — marquer les présences',
        '👨‍🎓 Étudiant — consulter son historique',
        '👪 Parent — voir les absences de son enfant',
    ], GREEN),
    ('⚙️ Stack Technique', [
        'Backend: Django 4.2 LTS',
        'Base de données: MySQL (XAMPP)',
        'Frontend: Bootstrap 5 + Jinja2',
        'Export: openpyxl, reportlab',
    ], ORANGE),
]
for i, (title, items, color) in enumerate(cols):
    x = 0.3 + i * 4.3
    add_rect(s, x, 2.9, 4.1, 4.2, fill_color=WHITE)
    add_rect(s, x, 2.9, 4.1, 0.5, fill_color=color)
    add_text(s, title, x + 0.15, 2.95, 3.9, 0.42, size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(s, '• ' + item, x + 0.15, 3.52 + j * 0.75, 3.85, 0.65,
                 size=12, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — ARCHITECTURE TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Architecture Technique', 'Structure du projet Django')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

apps = [
    ('accounts/',      'Utilisateurs, Filières, Groupes, Étudiants, Enseignants, Parents', BLUE_DARK),
    ('cours/',         'Cours, Sessions de cours, Emploi du temps', GREEN),
    ('presences/',     'Enregistrement et historique des présences', ORANGE),
    ('enseignants/',   'Vues spécifiques aux enseignants', RGBColor(0x6F,0x42,0xC1)),
    ('etudiants/',     'Liste et détail des étudiants', RGBColor(0x0D,0xCA,0xF0)),
    ('rapports/',      'Statistiques et exports Excel/PDF', RGBColor(0xDC,0x35,0x45)),
    ('notifications/', 'Alertes automatiques aux parents', RGBColor(0xFD,0x7E,0x14)),
]

add_text(s, 'Les 7 Applications Django', 0.4, 1.55, 8, 0.38, size=16, bold=True, color=BLUE_DARK)
for i, (app, desc, color) in enumerate(apps):
    y = 2.0 + i * 0.72
    add_rect(s, 0.4, y, 2.2, 0.62, fill_color=color)
    add_text(s, app, 0.5, y + 0.1, 2.0, 0.45, size=13, bold=True, color=WHITE)
    add_text(s, desc, 2.75, y + 0.12, 7.5, 0.42, size=12, color=DARK)

# Right: project tree
add_text(s, 'Structure', 10.6, 1.55, 2.7, 0.38, size=14, bold=True, color=BLUE_DARK)
tree_code = """gespresence/
├── accounts/
├── cours/
├── presences/
├── etudiants/
├── enseignants/
├── rapports/
├── notifications/
├── templates/
├── static/
│   ├── css/
│   └── js/
├── media/
├── .env
└── manage.py"""
code_box(s, tree_code, 10.5, 2.0, 2.7, 5.1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — MODÈLES DU PROJET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Modèles du Projet — MLD Simplifié', 'Relations entre les entités principales')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

entities = [
    (0.3,  1.6, 'Utilisateur',  ['id', 'username', 'email', 'role', 'telephone'],              BLUE_DARK),
    (3.3,  1.6, 'Etudiant',     ['id', 'user_id FK', 'matricule', 'groupe_id FK'],              GREEN),
    (6.3,  1.6, 'Enseignant',   ['id', 'user_id FK', 'specialite', 'cin'],                      ORANGE),
    (9.3,  1.6, 'Parent',       ['id', 'user_id FK', 'enfants M2M'],                             RGBColor(0x6F,0x42,0xC1)),
    (0.3,  4.5, 'Filiere',      ['id', 'nom', 'code'],                                           RGBColor(0x0D,0xCA,0xF0)),
    (3.3,  4.5, 'Groupe',       ['id', 'nom', 'annee', 'filiere_id FK'],                         RGBColor(0xDC,0x35,0x45)),
    (6.3,  4.5, 'SessionCours', ['id', 'date', 'heure', 'cours_id FK', 'groupe_id FK'],          RGBColor(0xFD,0x7E,0x14)),
    (9.3,  4.5, 'Presence',     ['id', 'statut', 'etudiant_id FK', 'session_id FK'],             BLUE_DARK),
]

for (x, y, name, fields, color) in entities:
    h = 0.4 + len(fields) * 0.32
    add_rect(s, x, y, 2.9, h, fill_color=WHITE)
    add_rect(s, x, y, 2.9, 0.38, fill_color=color)
    add_text(s, name, x + 0.1, y + 0.05, 2.7, 0.3, size=12, bold=True, color=WHITE)
    for j, field in enumerate(fields):
        bg = RGBColor(0xF0,0xF4,0xFF) if j % 2 == 0 else WHITE
        add_rect(s, x, y + 0.38 + j * 0.32, 2.9, 0.32, fill_color=bg)
        add_text(s, field, x + 0.1, y + 0.4 + j * 0.32, 2.7, 0.3, size=10, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — FONCTIONNALITÉS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Fonctionnalités par Rôle', 'Ce que chaque utilisateur peut faire')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

roles_features = [
    ('👤 Administrateur', BLUE_DARK, [
        'Dashboard avec statistiques temps réel',
        'CRUD: Étudiants, Enseignants, Cours',
        'Gestion Filières, Groupes',
        'Génération rapports + export Excel',
        'Accès Django Admin complet',
    ]),
    ('👨‍🏫 Enseignant', GREEN, [
        'Dashboard: Sessions du jour',
        'Feuille de présence interactive',
        'Boutons "Tout Présent / Tout Absent"',
        'Justification des absences',
        'Consultation historique',
    ]),
    ('👨‍🎓 Étudiant', ORANGE, [
        'Dashboard: Taux de présence',
        'Histogramme présents/absents',
        'Historique détaillé par cours',
        'Statuts: Présent/Absent/Retard/Justifié',
    ]),
    ('👪 Parent', RGBColor(0x6F,0x42,0xC1), [
        'Dashboard: Taux de l\'enfant',
        'Notifications automatiques',
        '→ SMS-like lors d\'absence',
        'Marquer notifications comme lues',
        'Historique des alertes',
    ]),
]

for i, (role, color, features) in enumerate(roles_features):
    x = 0.3 + i * 3.2
    h_total = 0.5 + len(features) * 0.58 + 0.3
    add_rect(s, x, 1.6, 3.0, h_total, fill_color=WHITE)
    add_rect(s, x, 1.6, 3.0, 0.5, fill_color=color)
    add_text(s, role, x + 0.1, 1.65, 2.8, 0.42, size=13, bold=True, color=WHITE)
    for j, feat in enumerate(features):
        add_text(s, '✓ ' + feat, x + 0.12, 2.22 + j * 0.58, 2.8, 0.52,
                 size=11, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — FEUILLE DE PRÉSENCE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Fonctionnalité Clé : Feuille de Présence', 'Le cœur du système')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

# Mock UI
add_rect(s, 0.3, 1.6, 8.5, 5.5, fill_color=WHITE)
add_rect(s, 0.3, 1.6, 8.5, 0.6, fill_color=BLUE_DARK)
add_text(s, '✅ Feuille de Présence — HTML/CSS & Bootstrap — DEV-101 — 07/07/2026',
         0.4, 1.67, 8.3, 0.46, size=11, bold=True, color=WHITE)
# Header row
add_rect(s, 0.3, 2.22, 8.5, 0.38, fill_color=RGBColor(0xF0,0xF4,0xFF))
for label, x, w in [('#', 0.3, 0.5), ('Nom & Prénom', 0.82, 2.8),
                     ('Matricule', 3.64, 1.5), ('Présent', 5.16, 0.9),
                     ('Absent', 6.08, 0.85), ('Retard', 6.95, 0.85), ('Justif.', 7.82, 0.98)]:
    add_text(s, label, x + 0.05, 2.27, w - 0.1, 0.3, size=10, bold=True, color=DARK)

students_mock = [
    ('01', 'Ahmed Benaissa',   'DEV-101-01', '●', '○', '○', ''),
    ('02', 'Sara Alami',       'DEV-101-02', '●', '○', '○', ''),
    ('03', 'Karim Zidane',     'DEV-101-03', '○', '●', '○', 'Malade'),
    ('04', 'Fatima Chraibi',   'DEV-101-04', '○', '○', '●', ''),
    ('05', 'Yassine Benali',   'DEV-101-05', '●', '○', '○', ''),
]
colors_row = [WHITE, RGBColor(0xF8,0xF9,0xFA)]
for i, (n, nom, mat, pr, ab, re, jus) in enumerate(students_mock):
    y = 2.62 + i * 0.42
    add_rect(s, 0.3, y, 8.5, 0.4, fill_color=colors_row[i % 2])
    for val, x, w, clr in [
        (n,   0.3,  0.5,  GRAY),
        (nom, 0.82, 2.8,  DARK),
        (mat, 3.64, 1.5,  GRAY),
        (pr,  5.16, 0.9,  GREEN if pr=='●' else GRAY),
        (ab,  6.08, 0.85, RGBColor(0xDC,0x35,0x45) if ab=='●' else GRAY),
        (re,  6.95, 0.85, ORANGE if re=='●' else GRAY),
        (jus, 7.82, 0.98, DARK),
    ]:
        add_text(s, val, x + 0.05, y + 0.08, w - 0.1, 0.3, size=11, color=clr)

# Buttons
for label, x, color in [('✓ Tout Présent', 0.35, GREEN),
                          ('✗ Tout Absent', 2.05, RGBColor(0xDC,0x35,0x45)),
                          ('💾 Enregistrer', 6.3, BLUE_DARK)]:
    add_rect(s, x, 4.92, 1.6, 0.42, fill_color=color)
    add_text(s, label, x + 0.08, 4.97, 1.5, 0.32, size=11, bold=True, color=WHITE)

# Right panel: code
add_text(s, 'Code — views.py', 9.0, 1.65, 4.2, 0.38, size=14, bold=True, color=BLUE_DARK)
code10 = """@login_required
def feuille_presence(request, pk):
  session = get_object_or_404(
      SessionCours, pk=pk)
  etudiants = Etudiant.objects.filter(
      groupe=session.groupe)

  if request.method == 'POST':
    for e in etudiants:
      statut = request.POST.get(
          f'statut_{e.pk}', 'absent')
      Presence.objects.update_or_create(
          etudiant=e,
          session=session,
          defaults={'statut': statut}
      )
      # Notif parent si absent
      if statut in ('absent','retard'):
          notifier_parent(e, session)
    return redirect('session_list', pk)

  return render(request,
      'presences/feuille.html', {
          'session':   session,
          'etudiants': etudiants,
      })"""
code_box(s, code10, 9.0, 2.08, 4.2, 5.1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — DASHBOARD ADMIN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'Dashboard Administrateur', 'Vue d\'ensemble en temps réel')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF0,0xF4,0xFF))

# Mock navbar
add_rect(s, 0, 1.4, 13.33, 0.5, fill_color=BLUE_DARK)
add_text(s, '🎓 GesPresence        Admin Système  [Admin]   Déconnexion',
         0.2, 1.46, 13, 0.38, size=11, color=WHITE)

# Sidebar mock
add_rect(s, 0, 1.9, 2.2, 5.6, fill_color=WHITE)
sidebar_items = ['📊 Dashboard ●', '👥 Étudiants', '👨‍🏫 Enseignants', '📚 Cours',
                 '✅ Présences', '📈 Rapports', '🔔 Notifications', '⚙️ Admin Django']
for i, item in enumerate(sidebar_items):
    bg = BLUE_LIGHT if i == 0 else WHITE
    add_rect(s, 0, 1.92 + i * 0.65, 2.18, 0.62, fill_color=bg)
    add_text(s, item, 0.15, 1.98 + i * 0.65, 2.0, 0.48, size=11,
             color=BLUE_DARK if i == 0 else DARK)

# Stats cards
stats = [
    ('35', 'Étudiants', BLUE_DARK),
    ('6',  'Enseignants', GREEN),
    ('15', 'Cours', ORANGE),
    ('0',  'Absences\naujourd\'hui', RGBColor(0xDC,0x35,0x45)),
]
for i, (val, label, color) in enumerate(stats):
    x = 2.35 + i * 2.65
    add_rect(s, x, 2.0, 2.5, 1.4, fill_color=WHITE)
    add_rect(s, x, 2.0, 0.7, 1.4, fill_color=color)
    add_text(s, val, x + 0.85, 2.12, 1.5, 0.7, size=28, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, label, x + 0.78, 2.85, 1.65, 0.45, size=10, color=GRAY,
             align=PP_ALIGN.CENTER)

# Sessions table mock
add_rect(s, 2.35, 3.55, 10.65, 3.9, fill_color=WHITE)
add_text(s, 'Dernières Sessions', 2.5, 3.62, 5, 0.38, size=13, bold=True, color=DARK)
add_text(s, 'Génerer un rapport →', 9.5, 3.62, 3.3, 0.38, size=11, color=BLUE_DARK)
add_rect(s, 2.35, 4.05, 10.65, 0.35, fill_color=RGBColor(0xF0,0xF4,0xFF))
for label, x in [('Date',0.1),('Cours',1.0),('Groupe',3.5),('Statut',6.0)]:
    add_text(s, label, 2.35 + x, 4.1, 1.5, 0.28, size=10, bold=True, color=GRAY)
sessions = [
    ('01/08/2026','Mathématiques pour Dev','DEV-101','● Planifiée'),
    ('01/08/2026','Mathématiques Réseaux', 'RES-101', '● Planifiée'),
    ('01/08/2026','Fiscalité Marocaine',   'COMPTA-101','● Planifiée'),
]
for i, (dt, cours, grp, stat) in enumerate(sessions):
    y = 4.45 + i * 0.9
    bg = WHITE if i % 2 == 0 else RGBColor(0xF8,0xF9,0xFA)
    add_rect(s, 2.35, y, 10.65, 0.85, fill_color=bg)
    for val, x in [(dt,0.1),(cours,1.0),(grp,3.5),(stat,6.0)]:
        add_text(s, val, 2.35 + x, y + 0.2, 2.3, 0.45, size=11, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — RÉSULTATS & CHIFFRES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, 'GesPresence — Résultats & Données', 'Ce qui a été réalisé')
add_rect(s, 0, 1.4, 13.33, 6.1, fill_color=RGBColor(0xF8,0xF9,0xFA))

numbers = [
    ('35',  'Étudiants créés',         '6 groupes, 4 filières',       BLUE_DARK),
    ('35',  'Comptes Parents',         'Un parent par étudiant',       GREEN),
    ('6',   'Professeurs',             'Avec spécialités + grades',    ORANGE),
    ('15',  'Cours / Modules',         'Assignés par filière',         RGBColor(0x6F,0x42,0xC1)),
    ('96',  'Sessions planifiées',     '4 semaines, emploi du temps',  RGBColor(0xDC,0x35,0x45)),
    ('7',   'Apps Django',             '+ 1 admin Django',             RGBColor(0x0D,0xCA,0xF0)),
    ('21',  'Tables MySQL',            'Migrations appliquées',        RGBColor(0xFD,0x7E,0x14)),
    ('14',  'Templates HTML',          'Interface style OFPPT',        BLUE_DARK),
]

for i, (num, label, desc, color) in enumerate(numbers):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.25
    y = 1.65 + row * 2.1
    add_rect(s, x, y, 3.05, 1.85, fill_color=WHITE)
    add_rect(s, x, y, 3.05, 0.45, fill_color=color)
    add_text(s, num, x + 0.1, y + 0.05, 1.2, 0.38, size=22, bold=True, color=WHITE)
    add_text(s, label, x + 0.1, y + 0.55, 2.85, 0.52, size=14, bold=True, color=DARK)
    add_text(s, desc, x + 0.1, y + 1.12, 2.85, 0.52, size=11, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=BLUE_DARK)
add_rect(s, 0, 5.5, 13.33, 2.0, fill_color=RGBColor(0x07,0x46,0x8A))
add_rect(s, 0, 5.45, 13.33, 0.12, fill_color=GREEN)

add_text(s, 'Ce que vous avez appris aujourd\'hui', 0.6, 0.7, 12, 0.6,
         size=20, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.LEFT)
add_text(s, 'Django Fullstack', 0.6, 1.35, 12, 0.9,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

learned = [
    '✅  Architecture MVT et flux d\'une requête HTTP',
    '✅  Modèles ORM, migrations, relations de base de données',
    '✅  Vues, URLs, décorateurs d\'authentification',
    '✅  Templates, filtres, héritage de layouts',
    '✅  Formulaires, validation, protection CSRF',
    '✅  Système d\'authentification multi-rôles',
    '✅  Django Admin automatique',
    '✅  Projet complet : GesPresence (35 étudiants, 6 profs, 96 sessions)',
]
for i, item in enumerate(learned):
    col = i // 4
    row = i % 4
    x = 0.6 + col * 6.3
    y = 2.5 + row * 0.72
    add_text(s, item, x, y, 6.0, 0.62, size=14, color=WHITE)

add_text(s, '🌐  http://127.0.0.1:8000/', 0.6, 6.0, 6, 0.5,
         size=16, bold=True, color=RGBColor(0x88,0xFF,0x99))
add_text(s, '👤  admin / admin1234', 6.5, 6.0, 5, 0.5,
         size=14, color=RGBColor(0xCC,0xDD,0xFF))
add_text(s, '2026 — GesPresence — Django 4.2 LTS', 0.6, 6.8, 12, 0.4,
         size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
OUTPUT = r"C:\Users\Abdeladim\Desktop\Django_GesPresence_Presentation.pptx"
prs.save(OUTPUT)
print(f"Presentation saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
